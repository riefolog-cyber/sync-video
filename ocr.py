#!/usr/bin/env python3
"""
FASE 1 — Estrazione slide da PDF + OCR parallelo.
Il rendering PDF usa multiprocessing (PyMuPDF non è thread-safe).
"""

import contextlib
import hashlib
import io
import os
import re
import shutil
import subprocess
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from multiprocessing import Pool
from pathlib import Path

import pymupdf as fitz  # PyMuPDF (alias deprecato importato come compat)
import pytesseract
from PIL import Image

from config import log, tqdm

# Candidati percorsi LibreOffice (Windows + Unix)
_SOFFICE_CANDIDATES = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/soffice",
    "/usr/local/bin/soffice",
    "/opt/libreoffice/program/soffice",
]

# Candidati percorsi ONLYOFFICE (x2t converter)
_ONLYOFFICE_CANDIDATES = [
    r"C:\Program Files\ONLYOFFICE\DesktopEditors\converter\x2t.exe",
    r"C:\Program Files (x86)\ONLYOFFICE\DesktopEditors\converter\x2t.exe",
    "/usr/bin/x2t",
    "/opt/onlyoffice/desktopeditors/converter/x2t",
]


def _find_soffice() -> str | None:
    """Cerca un eseguibile LibreOffice soffice utilizzabile per la conversione."""
    env = os.environ.get("SOFFICE_PATH")
    if env and Path(env).exists():
        return env
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    for cand in _SOFFICE_CANDIDATES:
        if Path(cand).exists():
            return cand
    return None


def _find_onlyoffice() -> str | None:
    """Cerca il converter ONLYOFFICE x2t utilizzabile per la conversione."""
    env = os.environ.get("ONLYOFFICE_X2T_PATH")
    if env and Path(env).exists():
        return env
    found = shutil.which("x2t")
    if found:
        return found
    for cand in _ONLYOFFICE_CANDIDATES:
        if Path(cand).exists():
            return cand
    return None


PRESENTATION_SUFFIXES = {".ppt", ".pptx"}


def _file_md5(path: Path, chunk: int = 1024 * 1024) -> str:
    """MD5 del contenuto di un file (streaming, non carica in memoria)."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(chunk), b""):
            h.update(block)
    return h.hexdigest()


def _pptx_fallback_to_pdf(ppt_path: Path, pdf_path: Path) -> bool:
    """Fallback puro Python: PPTX -> PDF via python-pptx + PyMuPDF.

    Estrae testo e immagini dalle slide e genera un PDF 16:9 (1280x720).
    Non è fedele al 100% come LibreOffice/ONLYOFFICE, ma sblocca la pipeline
    quando il converter nativo fallisce (es. x2t TypeError su file complessi).
    Ritorna True se il PDF è stato creato con successo.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return False
    try:
        prs = Presentation(str(ppt_path))
    except Exception as e:
        log.debug("   Fallback PPTX->PDF: apertura fallita: %s", e)
        return False
    try:
        doc = fitz.open()
        for idx, slide in enumerate(prs.slides):
            # Pagina 16:9 come il video finale (1280x720)
            page = doc.new_page(width=1280, height=720)
            # Raccogli testo
            texts: list[str] = []
            images_to_draw: list[tuple[bytes, str]] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
                # Immagini
                if shape.shape_type == 13:  # picture
                    try:
                        img_bytes = shape.image.blob
                        ext = shape.image.ext or "png"
                        images_to_draw.append((img_bytes, ext))
                    except Exception:
                        pass
            # Se ci sono immagini, prova a disegnarne la prima a tutta pagina
            if images_to_draw:
                try:
                    img_bytes, ext = images_to_draw[0]

                    # Fix RGBA -> RGB (fitz/x2t non gestisce alpha correttamente: rendeva bianco)
                    try:
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        if pil_img.mode in ("RGBA", "LA", "PA"):
                            bg = Image.new("RGB", pil_img.size, (255, 255, 255))
                            if pil_img.mode == "RGBA":
                                bg.paste(pil_img, mask=pil_img.split()[3])
                            else:
                                bg.paste(pil_img)
                            buf = io.BytesIO()
                            bg.save(buf, format="PNG")
                            img_bytes = buf.getvalue()
                            ext = "png"
                    except Exception:
                        pass

                    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tf:
                        tf.write(img_bytes)
                        tmp_img = tf.name
                    try:
                        page.insert_image(page.rect, filename=tmp_img, keep_proportion=True, overlay=False)
                    finally:
                        Path(tmp_img).unlink(missing_ok=True)
                    # Se c'era testo oltre l'immagine, aggiungilo in basso
                    if texts:
                        text = " | ".join(texts)[:300]
                        page.insert_textbox(
                            fitz.Rect(20, 680, 1260, 710),
                            text,
                            fontsize=8,
                            color=(0.3, 0.3, 0.3),
                            align=1,
                        )
                except Exception as e:
                    log.debug("   Fallback immagine slide %d fallita: %s", idx + 1, e)
            # Nessuna immagine: renderizza il testo centrato
            if not images_to_draw:
                full_text = "\n".join(texts).strip() or f"[Slide {idx+1} - contenuto visivo]"
                # Titolo slide
                rect = fitz.Rect(40, 40, 1240, 680)
                page.insert_textbox(
                    rect,
                    full_text[:2000],
                    fontsize=18,
                    color=(0.1, 0.1, 0.1),
                    align=0,
                    fontname="helv",
                )
            # Numero slide in basso a destra
            page.insert_text(fitz.Point(1220, 710), f"{idx+1}/{len(prs.slides)}", fontsize=8, color=(0.5, 0.5, 0.5))
        doc.save(str(pdf_path))
        doc.close()
        return pdf_path.exists() and pdf_path.stat().st_size > 0
    except Exception as e:
        log.debug("   Fallback PPTX->PDF eccezione: %s", e)
        with contextlib.suppress(Exception):
            doc.close()
        return False


def convert_presentation_to_pdf(ppt_path: Path, out_dir: Path) -> Path:
    """Converte una presentazione PPT/PPTX in PDF usando LibreOffice o ONLYOFFICE.

    Ordine di preferenza:
      1. LibreOffice ``soffice`` (se installato)
      2. ONLYOFFICE ``x2t`` (se installato) — alternativa senza LibreOffice
      3. Fallback Python puro (python-pptx + PyMuPDF) se i converter falliscono
    """
    soffice = _find_soffice()
    onlyoffice = _find_onlyoffice()
    has_converter = soffice or onlyoffice
    # Anche senza converter proviamo il fallback Python puro
    if not has_converter:
        log.warning("   Nessun converter nativo (LibreOffice/ONLYOFFICE) trovato: provo fallback Python puro...")
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = out_dir / (ppt_path.stem + ".pdf")

    # Riusa il PDF convertito solo se il file sorgente è INVARIATO (confronto per
    # contenuto, non per data di modifica): una copia/caricamento del .pptx può
    # avere un mtime più vecchio del PDF in cache, cosa che faceva riusare un PDF
    # stantio. Il marker a lato memorizza l'MD5 del sorgente convertito.
    marker_path = out_dir / (ppt_path.stem + ".src_md5")
    if pdf_path.exists() and marker_path.exists():
        try:
            if marker_path.read_text(encoding="ascii").strip() == _file_md5(ppt_path):
                log.info("   -> PDF già presente (riusato): %s", pdf_path)
                return pdf_path
        except OSError:
            pass
        log.info("   -> PDF in cache scaduto (sorgente cambiato), riconverto...")

    ext_upper = ppt_path.suffix.upper().lstrip(".")
    # Prova converter nativo se disponibile
    if has_converter:
        converter: str
        cmd: list[str]
        if soffice:
            converter = "LibreOffice"
            cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(ppt_path)]
        else:
            converter = "ONLYOFFICE"
            cmd = [onlyoffice, str(ppt_path), str(pdf_path)]  # type: ignore[list-item]
        log.info("   Conversione %s -> PDF (%s)...", ext_upper, converter)
        proc = None
        try:
            proc = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=180,
                check=False,
                creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0,
            )
        except subprocess.TimeoutExpired:
            log.warning("   Conversione %s->PDF timeout (180s), provo fallback...", ext_upper)
        if proc is not None:
            # x2t ritorna exit !=0 anche quando produce il PDF con warning:
            # il successo si valuta dalla presenza del file, non dall'exit code.
            if pdf_path.exists() and pdf_path.stat().st_size > 0:
                if proc.returncode != 0:
                    log.debug(
                        "   %s warning (exit=%d): %s",
                        converter,
                        proc.returncode,
                        (proc.stderr or proc.stdout)[:500],
                    )
                marker_path.write_text(_file_md5(ppt_path), encoding="ascii")
                log.info("   -> PDF convertito: %s", pdf_path)
                return pdf_path
            log.warning(
                "   %s fallito (exit=%d): %s",
                converter,
                proc.returncode,
                (proc.stderr or proc.stdout or "")[:600],
            )
            log.warning("   Provo fallback Python puro...")

    # Fallback puro Python (non richiede LibreOffice/ONLYOFFICE)
    log.info("   Conversione %s -> PDF (fallback Python python-pptx+PyMuPDF)...", ext_upper)
    if _pptx_fallback_to_pdf(ppt_path, pdf_path):
        marker_path.write_text(_file_md5(ppt_path), encoding="ascii")
        log.info("   -> PDF convertito (fallback): %s", pdf_path)
        log.warning(
            "   Nota: il fallback preserva testo/immagini ma non il layout esatto. "
            "Per fedeltà massima apri il PPTX in ONLYOFFICE e salva manualmente come PDF."
        )
        return pdf_path

    raise RuntimeError(
        f"Conversione {ext_upper}->PDF fallita con tutti i metodi.\n"
        f" - LibreOffice: {'non trovato' if not soffice else 'errore'}\n"
        f" - ONLYOFFICE x2t: {'non trovato' if not onlyoffice else 'errore (TypeError su questo file)'}\n"
        f" - Fallback Python: fallito\n"
        "Soluzioni immediate (senza installare LibreOffice):\n"
        " 1) Apri presentazione.pptx in ONLYOFFICE DesktopEditors > "
        "File > Scarica come > PDF e salva come presentazione.pdf\n"
        " 2) Oppure carica su Google Slides > File > Scarica > PDF\n"
        " 3) Poi lancia: "
        f"py -3.11 main.py --pdf \"{pdf_path}\""
    )


def _render_page(args: tuple[str, int, str, int]) -> str:
    """Renderizza una singola pagina PDF in PNG (eseguita in processo separato)."""
    pdf_path, page_num, output_path, dpi = args
    doc = fitz.open(str(pdf_path))
    try:
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=dpi)
        pix.save(str(output_path))
    finally:
        doc.close()
    return output_path


def _ocr_single_slide(image_path: Path, lang: str, max_retries: int = 3) -> str:
    """Esegue OCR su una singola immagine con retry e backoff."""
    raw = ""
    for attempt in range(max_retries):
        try:
            # Passiamo il PERCORSO, non l'oggetto Image: Pillow 12.x ha una
            # regressione quando pytesseract ri-salva il PNG (AttributeError:
            # 'PngImageFile' object has no attribute '_im'). Col path pytesseract
            # usa il file direttamente senza duplicarlo.
            raw = pytesseract.image_to_string(str(image_path), lang=lang)
            break
        except (pytesseract.TesseractError, OSError) as e:
            if attempt < max_retries - 1:
                wait = (attempt + 1) * 1.0
                log.debug(
                    "   OCR retry %d/%d per %s (errore: %s), attesa %.0fs...",
                    attempt + 1,
                    max_retries,
                    image_path.name,
                    e,
                    wait,
                )
                time.sleep(wait)
            else:
                log.debug("   OCR fallito con lang=%s, riprovo senza lingua.", lang)
                try:
                    raw = pytesseract.image_to_string(str(image_path))
                except (pytesseract.TesseractError, OSError):
                    raw = ""
    clean = re.sub(r"\s+", " ", raw).strip()
    return clean if clean else "[Nessun testo rilevato. Immagine visiva.]"


def extract_slides_text_ocr(
    pdf_path: Path,
    output_dir: Path,
    lang: str = "ita",
    dpi: int = 300,
    workers: int = 4,
) -> tuple[list[str], list[str]]:
    """Converte ogni pagina del PDF in PNG ed estrae il testo via OCR."""
    log.info("1. Analisi del PDF: Estrazione OCR dei testi per l'IA...")
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    total_pages = len(doc)
    doc.close()

    slide_files: list[str] = []
    pages_to_render: list[tuple[str, int, str, int]] = []

    for page_num in range(total_pages):
        output_path = output_dir / f"slide_{page_num:03d}.png"
        valid_cache = False
        if output_path.exists():
            try:
                with Image.open(str(output_path)) as test_img:
                    test_img.verify()
                valid_cache = True
            except (OSError, ValueError, SyntaxError):
                log.debug("   Cache corrotta per %s, ri-renderizzo.", output_path.name)
                output_path.unlink(missing_ok=True)
        if not valid_cache:
            pages_to_render.append((str(pdf_path), page_num, str(output_path), dpi))
        slide_files.append(str(output_path))

    if pages_to_render:
        n_render = len(pages_to_render)
        log.info("   Rendering %d slide (%d cached)...", n_render, total_pages - n_render)
        n_workers = min(workers, n_render)
        if n_workers <= 1 or n_render <= 1:
            for args in tqdm(pages_to_render, desc="Rendering slide"):
                _render_page(args)
        else:
            with Pool(processes=n_workers) as pool:
                list(tqdm(pool.imap(_render_page, pages_to_render), total=n_render, desc="Rendering slide"))
    else:
        log.info("   Tutte le slide sono in cache (nessun rendering necessario).")

    slide_texts: list[str] = [""] * total_pages
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {executor.submit(_ocr_single_slide, Path(sf), lang): idx for idx, sf in enumerate(slide_files)}
        for future in tqdm(as_completed(futures), total=total_pages, desc="OCR slide"):
            idx = futures[future]
            slide_texts[idx] = future.result()

    for i, txt in enumerate(slide_texts):
        preview = txt[:80] + "..." if len(txt) > 80 else txt
        log.debug("   Slide %d: %s", i + 1, preview)

    log.info("   -> OCR completato su %d slide.", total_pages)
    return slide_files, slide_texts
